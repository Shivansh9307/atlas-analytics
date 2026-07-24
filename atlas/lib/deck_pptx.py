"""python-pptx deck builder. Offline, no auth — the source-of-truth deck.

Fixed skeleton enforced here so every deck reads the same way:
  1. Title (a CLAIM, not a label)
  2. Key insight
  3..N Supporting evidence (one chart per claim)
  N+1  So-what
  N+2  Recommendation
  Appendix: Methodology, Assumptions, Provenance table

Contract: each content slide may reference provenance claim_ids. build_deck
collects them so GATE 4 can verify none is an orphan before export.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt


@dataclass
class Chart:
    kind: str                       # "bar" | "column" | "line"
    categories: list[str]
    series: dict[str, list[float]]  # series name -> values
    title: str = ""


@dataclass
class Slide:
    kind: str                       # insight|evidence|sowhat|recommendation
    title: str                      # a CLAIM, not a label
    bullets: list[str] = field(default_factory=list)
    chart: Chart | None = None
    speaker_notes: str = ""         # 60-90s spoken
    claim_ids: list[str] = field(default_factory=list)


@dataclass
class DeckSpec:
    title: str                      # headline claim
    subtitle: str
    decision_owner: str
    slides: list[Slide] = field(default_factory=list)
    assumptions: list[str] = field(default_factory=list)
    methodology: str = ""
    provenance: list[dict] = field(default_factory=list)  # {claim_id,text,value,query_hash,tier}

    def referenced_claim_ids(self) -> list[str]:
        ids: list[str] = []
        for s in self.slides:
            ids.extend(s.claim_ids)
        return ids


_CHART_TYPE = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
}


def build_deck(spec: DeckSpec, out_path: str | Path,
               notes_path: str | Path | None = None) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs, spec.title, f"{spec.subtitle}\nPrepared for: {spec.decision_owner}")

    notes_md = [f"# Speaker notes — {spec.title}\n"]
    for i, sl in enumerate(spec.slides, start=1):
        _content_slide(prs, sl)
        notes_md.append(f"## Slide {i + 1}: {sl.title}\n\n{sl.speaker_notes or '(no notes)'}\n")

    # --- appendix ---
    _bullets_slide(prs, "Appendix — Methodology",
                   [ln for ln in (spec.methodology or "").splitlines() if ln.strip()]
                   or ["(methodology not supplied)"])
    _bullets_slide(prs, "Appendix — Assumptions",
                   spec.assumptions or ["(no assumptions declared)"])
    _provenance_slide(prs, spec.provenance)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    if notes_path:
        Path(notes_path).write_text("\n".join(notes_md))
    return out_path


# --- slide builders ---
def _title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _content_slide(prs, sl: Slide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = sl.title

    if sl.bullets:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.6), Inches(5.3))
        tf = box.text_frame
        tf.word_wrap = True
        for j, b in enumerate(sl.bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(16)

    if sl.chart:
        _add_chart(slide, sl.chart,
                   left=Inches(6.6) if sl.bullets else Inches(1.5),
                   width=Inches(6.0) if sl.bullets else Inches(10.0))

    _set_notes(slide, sl.speaker_notes)


def _bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for j, b in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(15)


def _provenance_slide(prs, provenance: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Appendix — Provenance (every number is sourced)"
    rows = len(provenance) + 1
    cols = 4
    if rows == 1:
        _bullets_slide(prs, "Appendix — Provenance", ["(no provenance records)"])
        return
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5)
    ).table
    for c, head in enumerate(["Claim", "Value", "Query hash", "Evidence tier"]):
        table.cell(0, c).text = head
    for r, rec in enumerate(provenance, start=1):
        table.cell(r, 0).text = str(rec.get("text", ""))[:60]
        table.cell(r, 1).text = str(rec.get("value", ""))
        table.cell(r, 2).text = str(rec.get("query_hash", ""))
        table.cell(r, 3).text = str(rec.get("tier", ""))


def _add_chart(slide, chart: Chart, left, width):
    data = CategoryChartData()
    data.categories = chart.categories
    for name, vals in chart.series.items():
        data.add_series(name, vals)
    ct = _CHART_TYPE.get(chart.kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
    gframe = slide.shapes.add_chart(ct, left, Inches(1.4), width, Inches(5.2), data)
    ch = gframe.chart
    ch.has_legend = len(chart.series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    if chart.title:
        ch.has_title = True
        ch.chart_title.text_frame.text = chart.title
    else:
        ch.has_title = False


def _set_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text or ""
