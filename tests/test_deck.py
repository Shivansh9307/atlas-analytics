from pptx import Presentation

from atlas.lib.deck_pptx import Chart, DeckSpec, Slide, build_deck


def _spec():
    return DeckSpec(
        title="EMEA margin fell 4pts because volume shifted to low-margin Hardware",
        subtitle="Q2 vs Q1 gross-margin decomposition",
        decision_owner="VP Finance, EMEA",
        slides=[
            Slide(
                kind="insight",
                title="Gross margin dropped from 60.0% to 56.0% — a 4.0pt fall",
                bullets=["Software margin unchanged (70%)", "Hardware margin unchanged (20%)"],
                chart=Chart("column", ["Q1", "Q2"], {"Gross margin %": [60.0, 56.0]},
                            title="EMEA gross margin"),
                speaker_notes="The headline is a four point drop, entirely explained by mix.",
                claim_ids=["c_gm_q1", "c_gm_q2"],
            ),
            Slide(
                kind="evidence",
                title="The entire drop is a mix shift, not a rate collapse",
                bullets=["Mix effect: -4.0pts", "Rate effect: ~0", "Interaction: ~0"],
                chart=Chart("bar", ["Mix", "Rate", "Interaction"],
                            {"Contribution (pts)": [-4.0, 0.0, 0.0]}),
                speaker_notes="Decomposition attributes all of it to mix.",
                claim_ids=["c_mix", "c_rate"],
            ),
            Slide(kind="recommendation",
                  title="Protect margin by rebalancing the EMEA product mix",
                  bullets=["Incentivise Software attach", "Review Hardware discounting"],
                  speaker_notes="Two levers, both mix-facing."),
        ],
        assumptions=["'Margin' = gross margin = (revenue - cogs)/revenue",
                     "Comparison window: Q2 vs Q1, EMEA only"],
        methodology="Exact mix/rate/interaction decomposition at product_line grain.",
        provenance=[
            {"claim_id": "c_gm_q1", "text": "EMEA Q1 GM", "value": "60.0%",
             "query_hash": "abc123", "tier": "decomposed"},
            {"claim_id": "c_gm_q2", "text": "EMEA Q2 GM", "value": "56.0%",
             "query_hash": "def456", "tier": "decomposed"},
        ],
    )


def test_build_deck_structure(tmp_path):
    out = tmp_path / "deck.pptx"
    notes = tmp_path / "speaker_notes.md"
    build_deck(_spec(), out, notes)
    assert out.exists()

    prs = Presentation(str(out))
    # title + 3 content + methodology + assumptions + provenance = 7
    assert len(prs.slides) == 7

    # every content slide title is a claim (has a verb-ish sentence, not a bare label)
    content_titles = [prs.slides[i].shapes.title.text for i in (1, 2, 3)]
    assert all(len(t.split()) >= 4 for t in content_titles)

    # charts present on the two evidence slides
    def has_chart(slide):
        return any(sh.has_chart for sh in slide.shapes if sh.shape_type == 3 or hasattr(sh, "has_chart") and _safe_has_chart(sh))
    assert _count_charts(prs) >= 2

    # speaker notes exist on content slides
    assert prs.slides[1].notes_slide.notes_text_frame.text.strip() != ""

    # notes markdown written
    assert notes.exists() and "Speaker notes" in notes.read_text()


def test_referenced_claims_collected():
    spec = _spec()
    ids = spec.referenced_claim_ids()
    assert "c_gm_q1" in ids and "c_mix" in ids


def _safe_has_chart(sh):
    try:
        return sh.has_chart
    except Exception:
        return False


def _count_charts(prs):
    n = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if _safe_has_chart(sh):
                n += 1
    return n
